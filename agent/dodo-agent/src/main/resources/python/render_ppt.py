#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
✔ 支持新 JSON schema（type / content / url / fontLimit）
✔ 保留模板字体样式
✔ Group shape 支持
✔ image 字段：仅 URL 才替换，非 URL 保留模板图
✔ shape.name = JSON key 匹配
✔ 每页 background 支持（全屏背景图，最稳定）
✔ 自动清理模板页
✔ duplicate_slide 完整复制图片关系
"""

import argparse
import json
import os
import sys
from pptx import Presentation
from urllib.parse import urlparse
from io import BytesIO
import requests
from copy import deepcopy

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

# =========================
# 工具函数
# =========================

def is_url(text):
    if not isinstance(text, str):
        return False
    try:
        r = urlparse(text)
        return r.scheme and r.netloc
    except:
        return False


def download_image(url):
    try:
        r = requests.get(url, timeout=20)
        r.raise_for_status()
        return BytesIO(r.content)
    except Exception as e:
        print(f"❌ 图片下载失败: {url}, 错误: {e}")
        return None


# =========================
# 遍历所有 shape（支持 group）
# =========================

def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            for sub in iter_shapes(shape.shapes):
                yield sub


# =========================
# 保留样式替换文本（支持 fontLimit）
# =========================
def replace_text_keep_style(shape, value, font_limit=None):
    if not shape.has_text_frame:
        return

    text = "" if value is None else str(value)

    # --- 优化后的截断逻辑 ---
    if font_limit:
        # 如果长度超过限制 3 个字符及以上，执行截断
        # 如果只多出 1-2 个字符，则保留完整文本
        if len(text) > (font_limit + 2):
            text = text[:font_limit]
    # -----------------------

    tf = shape.text_frame

    if not tf.paragraphs:
        tf.text = text
        return

    p = tf.paragraphs[0]

    if not p.runs:
        p.add_run().text = text
    else:
        p.runs[0].text = text
        # 清理多余的 runs
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)

    # 清理多余的段落
    for extra in tf.paragraphs[1:]:
        extra._element.getparent().remove(extra._element)


# =========================
# 设置幻灯片背景
# =========================

def set_slide_background(slide, url):
    """
    设置 PPT 背景图（真正背景，不是shape）
    """

    img_data = download_image(url)
    if not img_data:
        return False

    try:
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import qn

        # 1️⃣ 添加图片到关系
        image_part, rId = slide.part.get_or_add_image_part(img_data)

        # 2️⃣ 删除旧背景
        slide_elm = slide._element
        for bg in slide_elm.findall(qn("p:bg")):
            slide_elm.remove(bg)

        # 3️⃣ 创建背景XML
        bg_xml = f"""
        <p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <p:bgPr>
            <a:blipFill dpi="0" rotWithShape="1">
              <a:blip r:embed="{rId}"/>
              <a:stretch>
                <a:fillRect/>
              </a:stretch>
            </a:blipFill>
          </p:bgPr>
        </p:bg>
        """

        bg_element = parse_xml(bg_xml)

        # 4️⃣ 插入背景
        slide_elm.insert(0, bg_element)

        print("✓ 背景图设置成功")
        return True

    except Exception as e:
        print("❌ 设置背景失败:", e)
        return False


# =========================
# ⭐ 复制幻灯片
# =========================

def duplicate_slide(prs, index):
    """
    复制幻灯片（精准抓取背景节点 + 完美复用媒体资源，杜绝文件异常膨胀）
    """
    source_slide = prs.slides[index]
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # ========= 1. 清空新页面自动生成的占位符 shape
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # ========= 2. 复制关系，建立旧 rId 到新 rId 的严格映射
    rId_mapping = {}
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        try:
            if rel.is_external:
                # 外部链接直接使用 target_ref
                new_rel = new_slide.part.rels.get_or_add_extRel(rel.reltype, rel.target_ref)
                rId_mapping[rel.rId] = new_rel.rId
            else:
                # 内部资源（图片等），必须通过 relate_to 绑定到原始的 target_part
                # 这能让新页面直接“引用”旧图片，避免文件体积翻倍，且关系绝对正确
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_mapping[rel.rId] = new_rId
        except Exception as e:
            pass

    def update_xml_rids(element):
        """遍历 XML 节点，把所有旧的 rId 替换为新的 rId"""
        for el in element.iter():
            for attr_name, attr_value in list(el.attrib.items()):
                if attr_value in rId_mapping:
                    el.set(attr_name, rId_mapping[attr_value])

    # ========= 3. 复制 shapes
    from copy import deepcopy
    for shape in source_slide.shapes:
        new_el = deepcopy(shape._element)
        update_xml_rids(new_el)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # ========= 4. 复制幻灯片独立背景（核心修复点）
    # 背景节点 <p:bg> 是嵌套在 <p:cSld> 下的！
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    source_bg = source_slide._element.cSld.find(f"{{{ns}}}bg")

    if source_bg is not None:
        # 清理新页面的旧背景
        new_bg = new_slide._element.cSld.find(f"{{{ns}}}bg")
        if new_bg is not None:
            new_slide._element.cSld.remove(new_bg)

        # 复制背景 XML 并更新里面的 rId
        bg_copy = deepcopy(source_bg)
        update_xml_rids(bg_copy)

        # 将背景插入到 <p:cSld> 的最前面
        new_slide._element.cSld.insert(0, bg_copy)

    return new_slide


# =========================
# 填充幻灯片（支持新 JSON schema）
# =========================

def fill_slide(slide, data):
    all_shapes = list(iter_shapes(slide.shapes))

    # name → shape
    shape_map = {
        (s.name or "").strip(): s
        for s in all_shapes if s.name
    }

    for key, field in data.items():

        if not isinstance(field, dict):
            continue

        field_type = field.get("type", "").lower()
        content = field.get("content")
        url = field.get("url")

        # 兼容 fontLimit / font-limit
        font_limit = field.get("fontLimit") or field.get("font-limit")

        # ========= background（无需 shape）
        if field_type == "background":
            if is_url(url):
                set_slide_background(slide, url)
            continue

        shape = shape_map.get(key)
        if not shape:
            print(f"⚠️ 未找到 shape: {key}")
            continue

        # ========= image
        if field_type == "image":
            if is_url(url):
                img_data = download_image(url)
                if img_data:
                    try:
                        left, top, width, height = (
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height
                        )

                        parent = shape._element.getparent()
                        parent.remove(shape._element)

                        slide.shapes.add_picture(
                            img_data,
                            left,
                            top,
                            width=width,
                            height=height
                        )

                        print(f"✓ 图片替换: {key}")

                    except Exception as e:
                        print(f"❌ 图片替换失败 {key}: {e}")

            # 非 URL 保留模板图
            continue

        # ========= text
        if field_type == "text":
            replace_text_keep_style(
                shape,
                content,
                font_limit
            )


# =========================
# 删除幻灯片
# =========================

def delete_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    prs.part.drop_rel(slides[index].rId)
    del slide_id_list[index]


# =========================
# 主渲染流程
# =========================

def render_ppt(template_path, schema_json, output_path):
    print("🚀 开始渲染流程...")

    try:
        schema = json.loads(schema_json)
    except Exception as e:
        print(f"❌ JSON 解析失败: {e}")
        return

    slides_data = schema.get("slides", [])
    prs = Presentation(template_path)
    template_count = len(prs.slides)

    print(f"📊 模板总页数: {template_count}")

    # 根据 JSON 生成新页面
    for i, slide_def in enumerate(slides_data):
        t_idx = slide_def.get("templatePageIndex", 1) - 1
        data = slide_def.get("data", {})

        if 0 <= t_idx < template_count:
            print(f"📝 渲染第 {i+1} 页 (模板 {t_idx+1})")
            new_slide = duplicate_slide(prs, t_idx)
            fill_slide(new_slide, data)
        else:
            print(f"⚠️ 模板索引越界: {t_idx+1}")

    # 删除模板页
    for i in reversed(range(template_count)):
        delete_slide(prs, i)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)

    print("✅ 渲染完成:", output_path)


# =========================
# CLI
# =========================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    schema_json = os.environ.get("PPT_SCHEMA")

    if not schema_json:
        schema_file = os.environ.get("PPT_SCHEMA_FILE")
        if schema_file and os.path.exists(schema_file):
            with open(schema_file, "r", encoding="utf-8") as f:
                schema_json = f.read()

    if not schema_json:
        print("❌ 未提供 PPT_SCHEMA 或 PPT_SCHEMA_FILE")
        sys.exit(1)

    render_ppt(args.template, schema_json, args.output)


if __name__ == "__main__":
    main()
